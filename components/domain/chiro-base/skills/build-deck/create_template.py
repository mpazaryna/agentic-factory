#!/usr/bin/env python3
"""
create_template.py — Generate the Chiro PPTX template with branded slide layouts.

Usage:
    python3 create_template.py [output.pptx]

Defaults to decks/templates/chiro-default.pptx

The template defines:
  - Color palette (navy, white, slate, green accent)
  - Font choices (Arial family)
  - Slide layouts with pre-positioned placeholders
  - Dark and light slide variants
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# --- Chiro palette ---
# These match the colors used in investor-deck-v1 and v2 Marp CSS

NAVY = RGBColor(0x0F, 0x17, 0x2A)       # dark background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_300 = RGBColor(0x94, 0xA3, 0xB8)   # muted body on dark
SLATE_600 = RGBColor(0x47, 0x55, 0x69)   # body text on light
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)   # headings on light
GREEN = RGBColor(0x10, 0xB9, 0x81)       # accent
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)    # off-white slide bg
CODE_BG = RGBColor(0xF1, 0xF5, 0xF9)     # code block background

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Fonts
FONT_HEADING = 'Arial'
FONT_BODY = 'Arial'
FONT_CODE = 'Courier New'


def create_template(output_path: str):
    """Build a minimal branded template PPTX.

    Since python-pptx can't create new slide layouts from scratch,
    we create a presentation with example slides that demonstrate
    each layout type. The marp_to_pptx.py script reads the palette
    from the template's first slide master's theme colors.

    More importantly, this template file can be opened in PowerPoint
    or Google Slides to edit the masters directly — giving the business
    team full control over the final look.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Slide 1: Title slide (dark) ---
    _add_title_slide(prs)

    # --- Slide 2: Section header (dark) ---
    _add_section_slide(prs)

    # --- Slide 3: Content slide (light) ---
    _add_content_slide(prs)

    # --- Slide 4: Content with code (light) ---
    _add_code_slide(prs)

    # --- Slide 5: Closing slide (dark) ---
    _add_closing_slide(prs)

    prs.save(output_path)
    print(f'Template saved to {output_path}')
    print(f'  {len(prs.slides)} example slides demonstrating layout patterns')
    print(f'  Edit slide masters in PowerPoint/Google Slides to customize')


def _set_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def _set_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def _add_accent_bar(slide):
    """Add a thin green accent line near the top."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(1.5),
        Inches(1.5), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


def _add_footer(slide, text='CHIRO', dark=True):
    """Add a subtle footer with branding."""
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.8),
        Inches(3), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = SLATE_300 if dark else SLATE_600
    run.font.name = FONT_BODY


def _add_title_slide(prs):
    """Dark title slide with centered text."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    _set_dark_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0),
        Inches(10.3), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'PRESENTATION TITLE'
    run.font.size = Pt(44)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_HEADING

    # Accent bar centered
    shape = slide.shapes.add_shape(
        1, Inches(5.9), Inches(3.3),
        Inches(1.5), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Subtitle
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.6),
        Inches(10.3), Inches(1.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Subtitle or tagline goes here'
    run.font.size = Pt(20)
    run.font.color.rgb = SLATE_300
    run.font.name = FONT_BODY

    _add_footer(slide, dark=True)


def _add_section_slide(prs):
    """Dark section divider slide."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_dark_bg(slide)

    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5),
        Inches(11.7), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Section Title'
    run.font.size = Pt(36)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_HEADING

    _add_accent_bar(slide)

    # Body
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.8),
        Inches(8), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Brief context or framing for this section of the presentation.'
    run.font.size = Pt(16)
    run.font.color.rgb = SLATE_300
    run.font.name = FONT_BODY

    _add_footer(slide, dark=True)


def _add_content_slide(prs):
    """Light content slide with title and bullets."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_light_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.7), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Slide Title'
    run.font.size = Pt(28)
    run.font.color.rgb = SLATE_800
    run.font.bold = True
    run.font.name = FONT_HEADING

    _add_accent_bar(slide)

    # Bullets
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8),
        Inches(11.7), Inches(4.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        'First point — keep bullets concise',
        'Second point — one idea per bullet',
        'Third point — use bold for emphasis',
    ]
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f'•  {text}'
        run.font.size = Pt(15)
        run.font.color.rgb = SLATE_600
        run.font.name = FONT_BODY

    _add_footer(slide, text='CHIRO', dark=False)


def _add_code_slide(prs):
    """Light slide with title and code block."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_light_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.7), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Code Example'
    run.font.size = Pt(28)
    run.font.color.rgb = SLATE_800
    run.font.bold = True
    run.font.name = FONT_HEADING

    _add_accent_bar(slide)

    # Code block
    code_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8),
        Inches(11.7), Inches(4.0)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG

    tf = code_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '# Code blocks render in monospace\ndef example():\n    return "editable text"'
    run.font.size = Pt(11)
    run.font.name = FONT_CODE
    run.font.color.rgb = SLATE_800

    _add_footer(slide, dark=False)


def _add_closing_slide(prs):
    """Dark closing slide."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_dark_bg(slide)

    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.2),
        Inches(10.3), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Thank You'
    run.font.size = Pt(40)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_HEADING

    # Accent bar
    shape = slide.shapes.add_shape(
        1, Inches(5.9), Inches(3.5),
        Inches(1.5), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Contact
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.9),
        Inches(10.3), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'contact@example.com'
    run.font.size = Pt(16)
    run.font.color.rgb = SLATE_300
    run.font.name = FONT_BODY

    _add_footer(slide, dark=True)


def main():
    output = sys.argv[1] if len(sys.argv) > 1 else None
    if not output:
        # Default location relative to repo root
        script_dir = Path(__file__).resolve().parent
        templates_dir = script_dir.parent / 'templates'
        templates_dir.mkdir(exist_ok=True)
        output = str(templates_dir / 'chiro-default.pptx')

    create_template(output)


if __name__ == '__main__':
    main()

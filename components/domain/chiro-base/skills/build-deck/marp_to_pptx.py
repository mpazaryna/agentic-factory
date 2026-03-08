#!/usr/bin/env python3
"""
marp_to_pptx.py — Convert Marp markdown presentations to editable PPTX.

Usage:
    python3 marp_to_pptx.py <input.md> [output.pptx]

If output path is omitted, writes to the same directory as input with .pptx extension.
Optionally accepts a --template path to a .pptx template file for consistent theming.
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# --- Slide model ---

def parse_marp(md_text: str) -> tuple[dict, list[dict]]:
    """Parse Marp markdown into frontmatter dict and list of slide dicts.

    Each slide dict has:
      - title: str or None
      - subtitle: str or None (## heading)
      - bullets: list[str]
      - code_blocks: list[dict] with keys 'lang' and 'code'
      - body: str (non-bullet, non-heading text)
      - image: str or None (first image path)
      - directives: dict (Marp directives from HTML comments)
    """
    # Split frontmatter
    frontmatter = {}
    body = md_text
    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', md_text, re.DOTALL)
    if fm_match:
        for line in fm_match.group(1).splitlines():
            if ':' in line:
                key, _, val = line.partition(':')
                frontmatter[key.strip()] = val.strip().strip('"\'')
        body = md_text[fm_match.end():]

    # Split into slides on --- (must be on its own line)
    raw_slides = re.split(r'\n---\s*\n', body)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slide = _parse_slide(raw)
        if slide:
            slides.append(slide)

    return frontmatter, slides


def _parse_slide(raw: str) -> dict | None:
    """Parse a single slide's raw markdown content."""
    lines = raw.splitlines()

    title = None
    subtitle = None
    bullets = []
    code_blocks = []
    body_lines = []
    image = None
    directives = {}

    in_code = False
    in_comment = False
    code_lang = ''
    code_lines = []

    for line in lines:
        # Marp directives in HTML comments
        dir_match = re.match(r'<!--\s*_(\w+):\s*(.+?)\s*-->', line)
        if dir_match:
            directives[dir_match.group(1)] = dir_match.group(2)
            continue

        # Skip HTML comments (single-line and multi-line)
        if re.match(r'^\s*<!--', line):
            if '-->' in line:
                continue
            # Start of multi-line comment
            in_comment = True
            continue
        if in_comment:
            if '-->' in line:
                in_comment = False
            continue

        # Code block fences
        if re.match(r'^```', line):
            if in_code:
                code_blocks.append({'lang': code_lang, 'code': '\n'.join(code_lines)})
                code_lines = []
                in_code = False
            else:
                in_code = True
                code_lang = line.lstrip('`').strip()
            continue

        if in_code:
            code_lines.append(line)
            continue

        # Headings (strip inline HTML tags)
        h1 = re.match(r'^#\s+(.+)', line)
        if h1 and title is None:
            title = re.sub(r'<[^>]+>', '', h1.group(1)).strip()
            continue

        h2 = re.match(r'^##\s+(.+)', line)
        if h2:
            clean_h2 = re.sub(r'<[^>]+>', '', h2.group(1)).strip()
            if subtitle is None:
                subtitle = clean_h2
            else:
                body_lines.append(f'**{clean_h2}**')
            continue

        # Images
        img = re.match(r'!\[.*?\]\((.+?)\)', line)
        if img and image is None:
            image = img.group(1)
            continue

        # Bullets (- or *)
        bullet = re.match(r'^(\s*)[-*]\s+(.+)', line)
        if bullet:
            indent = len(bullet.group(1))
            text = bullet.group(2).strip()
            level = 1 if indent >= 2 else 0
            bullets.append({'text': text, 'level': level})
            continue

        # Skip HTML tags and <br/>
        if re.match(r'^\s*<', line):
            continue

        # Markdown tables → convert to readable text
        if re.match(r'^\s*\|', line):
            # Skip alignment rows (|:---:|)
            if re.match(r'^\s*\|[\s:|-]+\|\s*$', line):
                continue
            # Strip pipes and clean cells
            cells = [c.strip() for c in line.strip('| \n').split('|')]
            cells = [c for c in cells if c]
            if cells:
                body_lines.append('  |  '.join(cells))
            continue

        # Everything else is body — strip inline HTML tags
        stripped = re.sub(r'<[^>]+>', '', line).strip()
        if stripped:
            body_lines.append(stripped)

    # Close unclosed code block
    if in_code and code_lines:
        code_blocks.append({'lang': code_lang, 'code': '\n'.join(code_lines)})

    body = '\n'.join(body_lines).strip()

    if not title and not bullets and not code_blocks and not body:
        return None

    return {
        'title': title,
        'subtitle': subtitle,
        'bullets': bullets,
        'code_blocks': code_blocks,
        'body': body,
        'image': image,
        'directives': directives,
    }


# --- PPTX generation ---

# Layout indices in default python-pptx template
LAYOUT_TITLE = 0       # Title Slide
LAYOUT_CONTENT = 1     # Title and Content
LAYOUT_BLANK = 6       # Blank

# Palette — matches decks/templates/chiro-default.pptx
COLOR_DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
COLOR_BODY = RGBColor(0x47, 0x55, 0x69)
COLOR_MUTED = RGBColor(0x94, 0xA3, 0xB8)
COLOR_ACCENT = RGBColor(0x10, 0xB9, 0x81)
COLOR_CODE_BG = RGBColor(0xF1, 0xF5, 0xF9)

# Fonts
FONT_HEADING = 'Arial'
FONT_BODY = 'Arial'
FONT_CODE = 'Courier New'


def _strip_md_formatting(text: str) -> list[tuple[str, bool]]:
    """Parse inline markdown bold into (text, is_bold) segments."""
    segments = []
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            segments.append((part[2:-2], True))
        else:
            if part:
                segments.append((part, False))
    return segments


def _add_text_with_formatting(tf, text: str, font_size: Pt, color: RGBColor,
                               bold: bool = False, alignment=None):
    """Add a paragraph with inline bold support."""
    p = tf.add_paragraph()
    if alignment:
        p.alignment = alignment

    segments = _strip_md_formatting(text)
    for i, (seg_text, seg_bold) in enumerate(segments):
        run = p.add_run()
        run.text = seg_text
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold or seg_bold

    return p


def build_pptx(frontmatter: dict, slides: list[dict],
               template_path: str | None = None,
               input_dir: Path | None = None) -> Presentation:
    """Build a PPTX Presentation from parsed slide data."""
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        is_title_slide = (i == 0 and slide_data.get('title')
                          and not slide_data.get('subtitle'))
        _add_slide(prs, slide_data, is_title_slide, input_dir)

    return prs


def _add_accent_bar(slide, left, top, dark=False):
    """Add a thin green accent bar."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top,
        Inches(1.5), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def _add_footer(slide, text='CHIRO', dark=True):
    """Add a subtle branded footer."""
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.8),
        Inches(3), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = COLOR_MUTED if dark else COLOR_BODY
    run.font.name = FONT_BODY


def _add_slide(prs: Presentation, data: dict, is_title: bool,
               input_dir: Path | None):
    """Add a single slide to the presentation."""
    layout = prs.slide_layouts[LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)

    is_dark = data.get('directives', {}).get('class') == 'dark'

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG if is_dark else COLOR_LIGHT_BG

    title_color = COLOR_WHITE if is_dark else COLOR_DARK_TEXT
    body_color = COLOR_MUTED if is_dark else COLOR_BODY

    left = Inches(0.8)
    top = Inches(0.6)
    width = Inches(11.7)
    current_top = top

    # Determine the visual heading — use title if present, fall back to subtitle
    heading = data.get('title') or data.get('subtitle')
    has_both = data.get('title') and data.get('subtitle')

    # Heading
    if heading:
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if is_title:
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(44) if is_title else Pt(28)
        run.font.color.rgb = title_color
        run.font.bold = True
        run.font.name = FONT_HEADING
        current_top += Inches(1.2) if is_title else Inches(0.9)

    # Accent bar
    if heading:
        if is_title:
            _add_accent_bar(slide, Inches(5.9), current_top - Inches(0.2))
        else:
            _add_accent_bar(slide, left, current_top - Inches(0.05))
        current_top += Inches(0.15)

    # Subtitle (only if title was also present — otherwise subtitle was used as heading)
    if has_both:
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = data['subtitle']
        run.font.size = Pt(22)
        run.font.color.rgb = title_color
        run.font.bold = True
        run.font.name = FONT_HEADING
        current_top += Inches(0.7)

    # Body text (standalone — no bullets or code)
    if data.get('body') and not data.get('bullets') and not data.get('code_blocks'):
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in data['body'].split('\n'):
            p = _add_text_with_formatting(tf, line, Pt(16), body_color,
                                           alignment=PP_ALIGN.CENTER if is_title else None)
            for run in p.runs:
                run.font.name = FONT_BODY
        current_top += Inches(2.0)

    # Bullets
    if data.get('bullets'):
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(data['bullets']):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()

            indent_level = bullet.get('level', 0)
            p.level = indent_level
            p.space_after = Pt(8)

            prefix = '  ' * indent_level + '•  '
            segments = _strip_md_formatting(bullet['text'])
            for k, (seg_text, seg_bold) in enumerate(segments):
                run = p.add_run()
                run.text = (prefix + seg_text) if k == 0 else seg_text
                run.font.size = Pt(15) if indent_level == 0 else Pt(13)
                run.font.color.rgb = body_color
                run.font.bold = seg_bold
                run.font.name = FONT_BODY

        current_top += Inches(0.45 * len(data['bullets']))

    # Body text (with bullets or code present)
    if data.get('body') and (data.get('bullets') or data.get('code_blocks')):
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in data['body'].split('\n'):
            p = _add_text_with_formatting(tf, line, Pt(14), body_color)
            for run in p.runs:
                run.font.name = FONT_BODY
        current_top += Inches(1.0)

    # Code blocks
    for cb in data.get('code_blocks', []):
        code_height = min(Inches(0.2 * (cb['code'].count('\n') + 2)), Inches(4.5))
        txBox = slide.shapes.add_textbox(left, current_top, width, code_height)
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = COLOR_CODE_BG

        tf = txBox.text_frame
        tf.word_wrap = False

        code_lines = cb['code'].splitlines()
        if len(code_lines) > 20:
            code_lines = code_lines[:18] + ['    ...']

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '\n'.join(code_lines)
        run.font.size = Pt(10)
        run.font.name = FONT_CODE
        run.font.color.rgb = COLOR_DARK_TEXT

        current_top += code_height + Inches(0.2)

    # Image
    if data.get('image') and input_dir:
        img_path = input_dir / data['image']
        if img_path.exists():
            try:
                slide.shapes.add_picture(
                    str(img_path), left, current_top,
                    width=Inches(8), height=Inches(4)
                )
            except Exception:
                pass

    # Footer
    _add_footer(slide, dark=is_dark)


# --- CLI ---

def main():
    parser = argparse.ArgumentParser(
        description='Convert Marp markdown to editable PPTX'
    )
    parser.add_argument('input', help='Path to Marp .md file')
    parser.add_argument('output', nargs='?', help='Output .pptx path (default: same dir)')
    parser.add_argument('--template', help='Path to .pptx template file')
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f'Error: {input_path} not found', file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = Path(args.output).resolve()
    else:
        output_path = input_path.with_suffix('.pptx')

    md_text = input_path.read_text(encoding='utf-8')
    frontmatter, slides = parse_marp(md_text)

    print(f'Parsed {len(slides)} slides from {input_path.name}')

    prs = build_pptx(frontmatter, slides, args.template, input_path.parent)
    prs.save(str(output_path))

    print(f'Saved to {output_path}')


if __name__ == '__main__':
    main()
